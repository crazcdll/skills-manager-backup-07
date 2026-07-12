"""Build the weekly value-report PPT from a plan JSON.

Usage:
    python build_pptx.py <plan.json> <imgs_dir> <output.pptx>
                         [--skeleton path/to/skeleton.pptx]
                         [--ribbon   path/to/ribbon.png]

Defaults look for ../assets/skeleton.pptx and ../assets/ribbon.png relative to this file.

Plan JSON shape:
{
  "cover_title":  "...",        # e.g. "严管理稳落实，重训练促成长"
  "cover_team":   "...",        # e.g. "第四值周组：牛晨蕊、赵成程、…"
  "cover_date":   "2026.5.5",
  "footer":       null,   # optional; set to a string (e.g. "昌平校区 · 第四值周组") to enable
  "sections": [
    {"name": "亮点工作", "num": "01", "slides": [
        {"title": "...", "body": "...", "images": ["image1.jpeg", ...]},
        {"title": "...", "body_blocks": [["sub-head", "explanation"], ...], "images": []},
        ...
    ]},
    ...
  ]
}

Each slide must have exactly one of:
  - "body": one paragraph subtitle (rendered above the images)
  - "body_blocks": list of [label, explanation] pairs (rendered as colored label rows;
    used when there are no images for the slide)

Image layout is chosen automatically: 1 image -> centered; 2 -> side-by-side; 3 -> 3-col.
Aspect ratios are always preserved (images are fit inside boxes, never stretched).

Each image is added once. The builder does NOT enforce uniqueness — the plan is
expected to list each image in exactly one slide. The extractor produces a plan that
already satisfies this; verify if you edit it manually.
"""
from __future__ import annotations
import argparse
import json
import os
import sys
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_SKELETON = os.path.normpath(os.path.join(HERE, '..', 'assets', 'skeleton.pptx'))
DEFAULT_RIBBON   = os.path.normpath(os.path.join(HERE, '..', 'assets', 'ribbon.png'))

# Slide canvas (16:9, EMU)
SLIDE_W = 12192000
SLIDE_H = 6858000

# Brand
BLUE  = RGBColor(0x4A, 0x90, 0xC2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0xAA, 0xAA, 0xAA)
FONT  = '微软雅黑'

# Skeleton slide indices (must match the bundled skeleton.pptx)
COVER, TOC, SEC1_DIV, SEC2_DIV, SEC3_DIV, SEC4_DIV, END = 0, 1, 2, 3, 4, 5, 6


# ---------- helpers ----------

def _set_run_fonts(run, font=FONT):
    """Force latin/ea/cs typeface on a run so Chinese text renders consistently."""
    rPr = run._r.get_or_add_rPr()
    for child in list(rPr):
        tag = child.tag.split('}', 1)[-1]
        if tag in ('latin', 'ea', 'cs'):
            rPr.remove(child)
    for tag in ('latin', 'ea', 'cs'):
        el = etree.SubElement(rPr, qn(f'a:{tag}'))
        el.set('typeface', font)
        el.set('charset', '0')


def _add_run(p, text, *, size, bold=False, color=DARK, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_run_fonts(r, font)
    return r


def _set_paragraph(tf, text, *, size, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                   line_spacing=1.3):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    _add_run(p, text, size=size, bold=bold, color=color)
    return p


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None,
                line_spacing=1.3, margins=(91440, 91440, 45720, 45720)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = (Emu(m) for m in margins)
    tf.vertical_anchor = anchor
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    tb.line.fill.background()
    _set_paragraph(tf, text, size=size, bold=bold, color=color, align=align,
                   line_spacing=line_spacing)
    return tb


def add_picture_in_box(slide, image_path, left, top, box_w, box_h):
    """Fit picture into the box preserving aspect ratio; center within the box."""
    iw, ih = Image.open(image_path).size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    px = int(left + (box_w - w) / 2)
    py = int(top + (box_h - h) / 2)
    return slide.shapes.add_picture(image_path, px, py, width=w, height=h)


# ---------- header / footer ----------

def add_header(slide, num, name, ribbon_path):
    """Left: blue bar with '<num>  <name>' in white. Right: ribbon decoration."""
    bar = slide.shapes.add_shape(1, 0, Emu(369570), Emu(9994900), Emu(645160))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Emu(228600); tf.margin_right = Emu(91440)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _add_run(p, f'  {num}  ', size=28, bold=True, color=WHITE)
    _add_run(p, name,         size=28, bold=True, color=WHITE)
    # Ribbon top-right
    slide.shapes.add_picture(ribbon_path, Emu(10643870), 0,
                             width=Emu(1533525), height=Emu(1383665))


def add_footer(slide, text):
    add_textbox(slide, Emu(0), SLIDE_H - Emu(350000), SLIDE_W, Emu(300000),
                text, size=10, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.0)


# ---------- content slide layouts ----------

def add_text_blocks_slide(slide, blocks, *, body_left, body_top, body_width):
    """Render `blocks` as colored label cards down the slide. blocks = [(label, body), ...]"""
    block_top = body_top + Emu(150000)
    block_h = Emu(1100000)
    block_gap = Emu(200000)
    usable_h = SLIDE_H - block_top - Emu(500000)
    n = len(blocks)
    total_h = n * block_h + (n - 1) * block_gap
    if total_h > usable_h and n:
        block_h = max(Emu(700000), (usable_h - (n - 1) * block_gap) // n)

    label_w = Emu(2400000)
    text_left = body_left + label_w + Emu(180000)
    text_w = body_width - label_w - Emu(180000)

    for i, (label, txt) in enumerate(blocks):
        top = block_top + i * (block_h + block_gap)
        tag = slide.shapes.add_shape(1, body_left, top, label_w, block_h)
        tag.fill.solid(); tag.fill.fore_color.rgb = BLUE
        tag.line.fill.background()
        tf = tag.text_frame
        tf.margin_left = Emu(140000); tf.margin_right = Emu(140000)
        tf.margin_top = Emu(50000); tf.margin_bottom = Emu(50000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        _set_paragraph(tf, label, size=22, bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, line_spacing=1.1)
        add_textbox(slide, text_left, top, text_w, block_h, txt,
                    size=18, color=DARK, align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)


def add_content_slide(prs, sec_num, sec_name, slide_data, ribbon_path, footer):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header(slide, sec_num, sec_name, ribbon_path)

    title = slide_data.get('title', '') or ''
    body = slide_data.get('body', '') or ''
    body_blocks = slide_data.get('body_blocks') or []
    images = slide_data.get('images') or []

    body_left = Emu(360000)
    body_width = Emu(11400000)
    title_top = Emu(1080000)

    add_textbox(slide, body_left, title_top, body_width, Emu(550000),
                title, size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT,
                line_spacing=1.0)

    body_top = title_top + Emu(600000)

    if not images:
        # Text-only slide. Prefer body_blocks for visual interest.
        if body_blocks:
            add_text_blocks_slide(slide, body_blocks,
                                  body_left=body_left, body_top=body_top,
                                  body_width=body_width)
        elif body:
            add_textbox(slide, body_left, body_top, body_width, Emu(4500000),
                        body, size=24, color=DARK, align=PP_ALIGN.LEFT,
                        line_spacing=1.6)
    else:
        if body:
            body_h = Emu(900000)
            add_textbox(slide, body_left, body_top, body_width, body_h,
                        body, size=18, color=DARK, align=PP_ALIGN.LEFT,
                        line_spacing=1.4)
            img_top = body_top + body_h + Emu(150000)
        else:
            img_top = body_top + Emu(100000)

        img_area_h = SLIDE_H - img_top - Emu(450000)
        n = len(images)
        if n == 1:
            box_w = Emu(7000000)
            box_left = (SLIDE_W - box_w) // 2
            add_picture_in_box(slide, images[0], box_left, img_top, box_w, img_area_h)
        elif n == 2:
            gap = Emu(200000)
            box_w = (Emu(11400000) - gap) // 2
            for i, img in enumerate(images):
                add_picture_in_box(slide, img, body_left + i * (box_w + gap),
                                   img_top, box_w, img_area_h)
        else:
            # 3 or more — show first 3 across
            gap = Emu(150000)
            n3 = min(3, n)
            box_w = (Emu(11400000) - (n3 - 1) * gap) // n3
            for i, img in enumerate(images[:n3]):
                add_picture_in_box(slide, img, body_left + i * (box_w + gap),
                                   img_top, box_w, img_area_h)

    if footer:
        add_footer(slide, footer)
    return slide


# ---------- skeleton text replacement ----------

def replace_in_slide(slide, mapping):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for needle, replacement in mapping.items():
                    if needle in run.text:
                        run.text = run.text.replace(needle, replacement)


# ---------- main build ----------

def build(plan: dict, imgs_dir: str, out_path: str,
          skeleton_path: str = DEFAULT_SKELETON,
          ribbon_path: str = DEFAULT_RIBBON):
    prs = Presentation(skeleton_path)
    sld_lst = prs.slides._sldIdLst
    orig_ids = list(sld_lst)
    if len(orig_ids) != 7:
        raise SystemExit(f'Skeleton must have 7 slides, found {len(orig_ids)}')

    # Replace cover placeholders
    replace_in_slide(prs.slides[COVER], {
        '__COVER_TITLE__': plan.get('cover_title', ''),
        '__COVER_TEAM__':  plan.get('cover_team', ''),
        '__COVER_DATE__':  plan.get('cover_date', ''),
    })

    # Footer defaults to None (no footer). Set plan["footer"] to a string to enable.
    footer = plan.get('footer', None)

    # Build & track new slides per section
    new_ids_by_section = {}
    for sec in plan['sections']:
        added = []
        for sd in sec['slides']:
            # Resolve image paths
            sd2 = dict(sd)
            sd2['images'] = [os.path.join(imgs_dir, name) for name in sd.get('images', [])]
            add_content_slide(prs, sec.get('num', '01'), sec['name'], sd2,
                              ribbon_path, footer)
            added.append(list(sld_lst)[-1])
        new_ids_by_section[sec['name']] = added

    # Reorder: cover, toc, sec1div, sec1 content, sec2div, sec2 content, ...
    div_for = {
        '亮点工作': orig_ids[SEC1_DIV],
        '常规工作': orig_ids[SEC2_DIV],
        '问题建议': orig_ids[SEC3_DIV],
        '值周反思': orig_ids[SEC4_DIV],
    }
    final = [orig_ids[COVER], orig_ids[TOC]]
    for sec in plan['sections']:
        final.append(div_for[sec['name']])
        final.extend(new_ids_by_section[sec['name']])
    final.append(orig_ids[END])

    for sid in list(sld_lst):
        sld_lst.remove(sid)
    for sid in final:
        sld_lst.append(sid)

    prs.save(out_path)
    print(f'Saved {out_path} ({len(final)} slides)')


if __name__ == '__main__':
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('plan_json')
    ap.add_argument('imgs_dir')
    ap.add_argument('out_pptx')
    ap.add_argument('--skeleton', default=DEFAULT_SKELETON)
    ap.add_argument('--ribbon',   default=DEFAULT_RIBBON)
    args = ap.parse_args()

    with open(args.plan_json, encoding='utf-8') as f:
        plan = json.load(f)
    build(plan, args.imgs_dir, args.out_pptx,
          skeleton_path=args.skeleton, ribbon_path=args.ribbon)
